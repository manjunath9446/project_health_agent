from __future__ import annotations

from collections import Counter
from datetime import datetime
from pathlib import Path

from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession

from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from src.models.project_analysis import ProjectAnalysis


class PresentationService:

    # ==========================================================
    # INITIALIZATION
    # ==========================================================

    def __init__(self, session: AsyncSession):

        self.session = session

        self.colors = {

            "navy": RGBColor(15, 35, 62),

            "blue": RGBColor(15, 98, 254),

            "green": RGBColor(36, 161, 72),

            "amber": RGBColor(241, 194, 27),

            "red": RGBColor(218, 30, 40),

            "light": RGBColor(245, 247, 250),

            "dark": RGBColor(55, 65, 81),

            "white": RGBColor(255, 255, 255),

            "grey": RGBColor(120, 120, 120),

        }

    # ==========================================================
    # LOAD PROJECTS
    # ==========================================================

    async def load_projects(self):

        result = await self.session.execute(

            select(ProjectAnalysis)

        )

        return result.scalars().all()

    # ==========================================================
    # PRESENTATION
    # ==========================================================

    def create_presentation(self):

        prs = Presentation()

        prs.slide_width = Inches(13.333)

        prs.slide_height = Inches(7.5)

        return prs

    # ==========================================================
    # PORTFOLIO SUMMARY
    # ==========================================================

    async def portfolio_summary(self):

        projects = await self.load_projects()

        if not projects:

            raise RuntimeError(
                "No analysed projects found."
            )

        rag = Counter()

        total_health = 0

        total_delay = 0

        recommendations = []

        summaries = []

        risks = []

        for project in projects:

            rag[project.rag_status] += 1

            total_health += project.health_score

            total_delay += project.forecast_delay

            summaries.append(project.executive_summary)

            recommendations.append(project.recommendations)

            risks.append(project.risks)

        return {

            "projects": projects,

            "count": len(projects),

            "avg_health": round(
                total_health / len(projects),
                1,
            ),

            "avg_delay": round(
                total_delay / len(projects),
                1,
            ),

            "rag": rag,

            "recommendations": recommendations,

            "risks": risks,

            "summaries": summaries,

        }

    # ==========================================================
    # HELPER DRAWING
    # ==========================================================

    def add_title(self, slide, text):

        box = slide.shapes.add_textbox(

            Inches(0.5),

            Inches(0.3),

            Inches(10),

            Inches(0.7),

        )

        p = box.text_frame.paragraphs[0]

        p.text = text

        p.font.size = Pt(24)

        p.font.bold = True

        p.font.color.rgb = self.colors["navy"]

    def add_kpi_card(self, slide, left, top, title, value, color):

        shape = slide.shapes.add_shape(

            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,

            left,

            top,

            Inches(2.3),

            Inches(1.1),

        )

        fill = shape.fill

        fill.solid()

        fill.fore_color.rgb = self.colors["white"]

        line = shape.line

        line.color.rgb = color

        line.width = Pt(1.5)

        title_box = slide.shapes.add_textbox(

            left + Inches(0.08),

            top + Inches(0.08),

            Inches(2.1),

            Inches(0.35),

        )

        p = title_box.text_frame.paragraphs[0]

        p.text = title

        p.font.size = Pt(11)

        p.font.bold = True

        p.font.color.rgb = self.colors["dark"]

        value_box = slide.shapes.add_textbox(

            left + Inches(0.08),

            top + Inches(0.48),

            Inches(2.1),

            Inches(0.35),

        )

        p = value_box.text_frame.paragraphs[0]

        p.text = value

        p.font.size = Pt(16)

        p.font.bold = True

        p.font.color.rgb = color

    def add_text(self, slide, left, top, width, height, text, size=18):

        box = slide.shapes.add_textbox(left, top, width, height)

        tf = box.text_frame

        tf.clear()

        p = tf.paragraphs[0]

        p.text = text

        p.font.size = Pt(size)

        p.font.color.rgb = self.colors["dark"]

        p.alignment = PP_ALIGN.LEFT

    def add_footer(self, slide, page_no):

        footer = slide.shapes.add_textbox(

            Inches(0.5),

            Inches(7.0),

            Inches(4),

            Inches(0.2),

        )

        p = footer.text_frame.paragraphs[0]

        p.text = f"Page {page_no}"

        p.font.size = Pt(10)

        p.font.color.rgb = self.colors["grey"]

    # ==========================================================
    # COVER SLIDE
    # ==========================================================

    def cover_slide(self, prs):

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        background = slide.background.fill

        background.solid()

        background.fore_color.rgb = self.colors["navy"]

        title = slide.shapes.add_textbox(

            Inches(0.8),

            Inches(1.2),

            Inches(10),

            Inches(1),

        )

        p = title.text_frame.paragraphs[0]

        p.text = "Monthly Executive Portfolio Review"

        p.font.size = Pt(34)

        p.font.bold = True

        p.font.color.rgb = self.colors["white"]

        sub = slide.shapes.add_textbox(

            Inches(0.8),

            Inches(2.2),

            Inches(8),

            Inches(0.8),

        )

        p = sub.text_frame.paragraphs[0]

        p.text = "AI Project Health Intelligence"

        p.font.size = Pt(20)

        p.font.color.rgb = self.colors["white"]

        date = slide.shapes.add_textbox(

            Inches(0.8),

            Inches(6.2),

            Inches(4),

            Inches(0.4),

        )

        p = date.text_frame.paragraphs[0]

        p.text = datetime.now().strftime("%B %Y")

        p.font.size = Pt(18)

        p.font.color.rgb = self.colors["white"]

        return slide

    # ==========================================================
    # EXECUTIVE DASHBOARD
    # ==========================================================

    def executive_dashboard(self, prs, portfolio):

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        self.add_title(slide, "Executive Dashboard")

        self.add_kpi_card(

            slide,

            Inches(0.5),

            Inches(1.0),

            "Projects",

            str(portfolio["count"]),

            self.colors["blue"],

        )

        self.add_kpi_card(

            slide,

            Inches(3.2),

            Inches(1.0),

            "Health",

            f'{portfolio["avg_health"]}%',

            self.colors["green"]
            if portfolio["avg_health"] >= 80
            else self.colors["amber"]
            if portfolio["avg_health"] >= 60
            else self.colors["red"],

        )

        self.add_kpi_card(

            slide,

            Inches(5.9),

            Inches(1.0),

            "Forecast Delay",

            f'{portfolio["avg_delay"]} Days',

            self.colors["amber"],

        )

        self.add_kpi_card(

            slide,

            Inches(8.6),

            Inches(1.0),

            "RED Projects",

            str(portfolio["rag"].get("RED", 0)),

            self.colors["red"],

        )

        summary = f"""Portfolio Overview

• Projects Analysed : {portfolio['count']}

• Average Health : {portfolio['avg_health']}%

• Average Delay : {portfolio['avg_delay']} Days

Executive Observation

The AI assessment indicates that the current
portfolio requires active executive monitoring.

Projects classified as RED should be reviewed
immediately to reduce schedule slippage and
business delivery risk.
"""

        self.add_text(

            slide,

            Inches(0.6),

            Inches(2.6),

            Inches(12),

            Inches(3.5),

            summary,

            18,

        )

        self.add_footer(slide, 2)

        return slide

    # ==========================================================
    # PORTFOLIO INSIGHTS
    # ==========================================================

    def portfolio_insights(self, prs, portfolio):

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        self.add_title(slide, "Portfolio Insights")

        rag = portfolio["rag"]

        self.add_kpi_card(slide, Inches(0.6), Inches(1.0), "GREEN", str(rag.get("GREEN", 0)), self.colors["green"])

        self.add_kpi_card(slide, Inches(3.2), Inches(1.0), "AMBER", str(rag.get("AMBER", 0)), self.colors["amber"])

        self.add_kpi_card(slide, Inches(5.8), Inches(1.0), "RED", str(rag.get("RED", 0)), self.colors["red"])

        insights = f"""Executive Insights

• {rag.get('RED', 0)} projects require immediate leadership attention.

• Portfolio health is currently {portfolio['avg_health']}%.

• Forecast delay averages {portfolio['avg_delay']} days.

• AI detected recurring schedule-related risks.

• Recovery planning should focus on RED projects first.

Business Impact

Without intervention, schedule slippage is likely to increase across the portfolio.
"""

        self.add_text(slide, Inches(0.6), Inches(2.5), Inches(12), Inches(3.8), insights, 18)

        self.add_footer(slide, 3)

        return slide

    # ==========================================================
    # RISK DASHBOARD
    # ==========================================================

    def risk_dashboard(self, prs, portfolio):

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        self.add_title(slide, "Risk Intelligence")

        left = 0.7

        top = 1.0

        for i, project in enumerate(portfolio["projects"][:5]):

            card = slide.shapes.add_shape(

                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,

                Inches(left),

                Inches(top),

                Inches(5.8),

                Inches(0.75),

            )

            card.fill.solid()

            card.fill.fore_color.rgb = self.colors["light"]

            card.line.color.rgb = self.colors["blue"]

            txt = slide.shapes.add_textbox(

                Inches(left + 0.15),

                Inches(top + 0.08),

                Inches(5.3),

                Inches(0.55),

            )

            p = txt.text_frame.paragraphs[0]

            p.text = f"""

Project {project.project_id}

RAG : {project.rag_status}

Health : {project.health_score:.1f}%

Delay : {project.forecast_delay} Days

"""

            p.font.size = Pt(15)

            top += 0.9

        self.add_footer(slide, 4)

        return slide

    # ==========================================================
    # EXECUTIVE RECOMMENDATIONS
    # ==========================================================

    def executive_recommendations(self, prs):

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        self.add_title(slide, "Executive Recommendations")

        recommendations = [

            ("Recover Critical Path", self.colors["red"]),

            ("Increase Resource Allocation", self.colors["amber"]),

            ("Weekly Executive Reviews", self.colors["blue"]),

            ("Monitor High-Risk Projects", self.colors["green"]),

        ]

        x = 0.6

        y = 1.1

        for title, color in recommendations:

            card = slide.shapes.add_shape(

                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,

                Inches(x),

                Inches(y),

                Inches(5.9),

                Inches(0.9),

            )

            card.fill.solid()

            card.fill.fore_color.rgb = color

            card.line.color.rgb = color

            box = slide.shapes.add_textbox(

                Inches(x + 0.2),

                Inches(y + 0.18),

                Inches(5.4),

                Inches(0.4),

            )

            p = box.text_frame.paragraphs[0]

            p.text = title

            p.font.size = Pt(18)

            p.font.bold = True

            p.font.color.rgb = self.colors["white"]

            y += 1.05

        self.add_footer(slide, 5)

        return slide

    # ==========================================================
    # EXECUTIVE ROADMAP
    # ==========================================================

    def roadmap_slide(self, prs):

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        self.add_title(slide, "Executive Roadmap")

        steps = [

            ("Week 1", "Recover RED Projects", self.colors["red"]),

            ("Week 2", "Risk Review", self.colors["amber"]),

            ("Week 3", "Resource Reallocation", self.colors["blue"]),

            ("Week 4", "Portfolio Reassessment", self.colors["green"]),

        ]

        x = 0.7

        for week, title, color in steps:

            shape = slide.shapes.add_shape(

                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,

                Inches(x),

                Inches(2.3),

                Inches(2.3),

                Inches(1.2),

            )

            shape.fill.solid()

            shape.fill.fore_color.rgb = color

            shape.line.color.rgb = color

            box = slide.shapes.add_textbox(

                Inches(x + 0.15),

                Inches(2.45),

                Inches(2),

                Inches(0.8),

            )

            tf = box.text_frame

            p = tf.paragraphs[0]

            p.text = week

            p.font.bold = True

            p.font.size = Pt(18)

            p.font.color.rgb = self.colors["white"]

            p = tf.add_paragraph()

            p.text = title

            p.font.size = Pt(14)

            p.font.color.rgb = self.colors["white"]

            x += 2.9

        self.add_footer(slide, 6)

        return slide

    # ==========================================================
    # CLOSING SLIDE
    # ==========================================================

    def closing_slide(self, prs):

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        background = slide.background.fill

        background.solid()

        background.fore_color.rgb = self.colors["navy"]

        title = slide.shapes.add_textbox(

            Inches(1),

            Inches(2),

            Inches(10),

            Inches(1),

        )

        p = title.text_frame.paragraphs[0]

        p.text = "Thank You"

        p.font.size = Pt(36)

        p.font.bold = True

        p.font.color.rgb = self.colors["white"]

        sub = slide.shapes.add_textbox(

            Inches(1),

            Inches(3),

            Inches(9),

            Inches(1),

        )

        p = sub.text_frame.paragraphs[0]

        p.text = "AI Project Health Intelligence\nExecutive Portfolio Review"

        p.font.size = Pt(22)

        p.font.color.rgb = self.colors["white"]

        return slide

    # ==========================================================
    # GENERATE PRESENTATION
    # ==========================================================

    async def generate(self, output_folder="reports"):

        portfolio = await self.portfolio_summary()

        prs = self.create_presentation()

        self.cover_slide(prs)

        self.executive_dashboard(prs, portfolio)

        self.portfolio_insights(prs, portfolio)

        self.risk_dashboard(prs, portfolio)

        self.executive_recommendations(prs)

        self.roadmap_slide(prs)

        self.closing_slide(prs)

        output_dir = Path(output_folder)

        output_dir.mkdir(exist_ok=True, parents=True)

        filename = output_dir / f"Executive_Portfolio_{datetime.now():%Y%m%d_%H%M%S}.pptx"

        prs.save(filename)

        return str(filename)